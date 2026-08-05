"""
AutoPresent - PowerPoint Auto Presenter with Text-to-Speech
Reads slide notes aloud and automatically advances to the next slide.
"""

import tkinter as tk
from tkinter import filedialog, ttk, messagebox
import threading
import time
import os

from pptx import Presentation
from lxml import etree
import win32com.client
import pythoncom
import re

def split_into_sentences(text: str) -> list[str]:
    """Split text into sentences. Simple but effective."""
    if not text or not text.strip():
        return []
    
    # Split on . ! ? followed by space or end of string
    sentences = re.split(r'(?<=[.!?])\s+', text.strip())
    
    # Remove empty items
    return [s.strip() for s in sentences if s.strip()]

# XML namespaces used in PPTX notes slides
_PPTX_NSMAP = {
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
}
def split_into_chunks(text: str, max_words: int = 12) -> list[str]:
    """Split text into reasonably sized speaking chunks."""
    import re
    # Split on sentence endings first
    sentences = re.split(r'(?<=[.!?])\s+', text.strip())
    chunks = []
    for sent in sentences:
        words = sent.split()
        if not words:
            continue
        for i in range(0, len(words), max_words):
            chunk = " ".join(words[i:i + max_words])
            if chunk:
                chunks.append(chunk)
    return chunks

# ---------------------------------------------------------------------------
# TTS Engine wrapper (lives on a dedicated thread)
# ---------------------------------------------------------------------------

class TTSEngine:
    """
    Text-to-speech using Windows SAPI (SpVoice) directly via win32com.

    All speak() / stop() calls must come from the presenter background thread
    (the same thread that called CoInitialize).  The GUI thread only calls
    get_voices_sync(), set_rate(), set_volume(), set_voice() which are safe
    because they only read/write plain Python values protected by a lock.
    """

    # SAPI speak flags
    _SVSFDefault        = 0
    _SVSFlagsAsync      = 1   # non-blocking (we use blocking mode = 0)
    _SVSFPurgeBeforeSpeak = 2

    def __init__(self):
        self._voice = None          # win32com SpVoice — created on presenter thread
        self._lock  = threading.Lock()

        # Settings queued from GUI thread, applied before each speak()
        self._pending_rate:  int   | None = None   # SAPI rate: -10 .. +10
        self._pending_volume: int  | None = None   # SAPI volume: 0..100
        self._pending_voice_name: str | None = None

    # ------------------------------------------------------------------
    # GUI-thread helpers
    # ------------------------------------------------------------------    
    def get_voices_sync(self):
        """
        Return a list of objects with .name and .id attributes.
        Creates a temporary SpVoice just for enumeration, then releases it.
        """
        pythoncom.CoInitialize()
        try:
            v = win32com.client.Dispatch("SAPI.SpVoice")
            tokens = v.GetVoices()
            voices = []
            for i in range(tokens.Count):
                tok = tokens.Item(i)
                class _V:
                    pass
                vobj = _V()
                vobj.name = tok.GetDescription()
                vobj.id   = tok.GetDescription()   # use name as ID
                voices.append(vobj)
            del v
            return voices
        finally:
            pythoncom.CoUninitialize()

    def set_rate(self, rate_wpm: int):
        """Map words-per-minute (80-300) to SAPI rate (-10 to +10)."""
        # 175 wpm → 0, range roughly 80→-5, 300→+10
        sapi_rate = int((rate_wpm - 175) / 12.5)
        sapi_rate = max(-10, min(10, sapi_rate))
        with self._lock:
            self._pending_rate = sapi_rate

    def set_volume(self, volume: float):
        """Map 0.0-1.0 to SAPI volume 0-100."""
        with self._lock:
            self._pending_volume = int(volume * 100)

    def set_voice(self, voice_name: str):
        with self._lock:
            self._pending_voice_name = voice_name

    # ------------------------------------------------------------------
    # Presenter-thread methods
    # ------------------------------------------------------------------

    def _init_engine(self):
        """Create SpVoice on the presenter thread. Called once from _run()."""
        self._voice = win32com.client.Dispatch("SAPI.SpVoice")

    def _apply_pending(self):
        """Flush queued settings into the SpVoice object."""
        with self._lock:
            if self._pending_rate is not None:
                self._voice.Rate = self._pending_rate
                self._pending_rate = None
            if self._pending_volume is not None:
                self._voice.Volume = self._pending_volume
                self._pending_volume = None
            if self._pending_voice_name is not None:
                name = self._pending_voice_name
                self._pending_voice_name = None
                tokens = self._voice.GetVoices()
                for i in range(tokens.Count):
                    tok = tokens.Item(i)
                    if tok.GetDescription() == name:
                        self._voice.Voice = tok
                        break

    def speak_async(self, text: str):
        """Start speaking asynchronously. Must be called from the presenter thread."""
        self._apply_pending()
        # 1 = SVSFlagsAsync
        self._voice.Speak(text, 1)

    def is_speaking(self) -> bool:
        if not self._voice:
            return False
        try:
            # 2 = SRSEIsSpeaking
            return self._voice.Status.RunningState == 2
        except Exception:
            return False

    def stop(self):
        """
        Interrupt current speech.
        MUST be called from the presenter thread (the one that owns the SpVoice).
        """
        if self._voice:
            try:
                # Purge + async empty speak
                self._voice.Speak("", 1 | 2)   # SVSFlagsAsync | SVSFPurgeBeforeSpeak
            except Exception:
                pass
import wave
import winsound
import tempfile
from piper import PiperVoice

class PiperTTSEngine:
    def __init__(self):
        self._lock = threading.Lock()
        self._stop_flag = False
        self._is_playing = False
        self._pending_rate = 1.0
        self._pending_volume = 1.0
        self._current_voice_id = "en_US-lessac-high"
        self._voices = {}

        search_dirs = [
            os.path.expanduser("~/.local/share/piper/voices"),
            os.path.expanduser("~/AppData/Local/piper/voices"),
            os.path.dirname(__file__),
            ".",
        ]

        voice_files = {
            "en_US-lessac-high": "en_US-lessac-high.onnx",
            "en_US-amy-medium": "en_US-amy-medium.onnx",
        }

        for vid, filename in voice_files.items():
            for d in search_dirs:
                path = os.path.join(d, filename)
                if os.path.exists(path):
                    try:
                        self._voices[vid] = PiperVoice.load(path)
                        print(f"Loaded Piper voice: {vid}")
                        break
                    except Exception as e:
                        print(f"Failed to load {vid}: {e}")

        if not self._voices:
            raise FileNotFoundError("No Piper voice models found.")

        if "en_US-lessac-high" in self._voices:
            self._current_voice_id = "en_US-lessac-high"
        else:
            self._current_voice_id = next(iter(self._voices))

    def get_voices_sync(self):
        result = []
        display_names = {
            "en_US-lessac-high": "Piper - Lessac High (Neural)",
            "en_US-amy-medium": "Piper - Amy Medium (Neural)",
        }
        for vid in self._voices:
            class _V: pass
            v = _V()
            v.id = vid
            v.name = display_names.get(vid, vid)
            result.append(v)
        return result

    def set_rate(self, rate_wpm: int):
        scale = 175 / max(80, min(300, rate_wpm))
        with self._lock:
            self._pending_rate = scale

    def set_voice(self, voice_name: str):
        if voice_name in self._voices:
            with self._lock:
                self._current_voice_id = voice_name
                
    def set_volume(self, volume: float):
        with self._lock:
            self._pending_volume = max(0.0, min(1.0, float(volume)))

    def speak_async(self, text: str):
        self._stop_flag = False
        self._is_playing = True

        def _run():
            temp_path = None
            try:
                with self._lock:
                    length_scale = getattr(self, "_pending_rate", 1.0)
                    volume = getattr(self, "_pending_volume", 1.0)
                    voice = self._voices[self._current_voice_id]

                with tempfile.NamedTemporaryFile(delete=False, suffix=".wav") as f:
                    temp_path = f.name

                with wave.open(temp_path, "wb") as wav_file:
                    # Preferred modern way
                    try:
                        from piper.config import SynthesisConfig
                        syn_config = SynthesisConfig(length_scale=length_scale)
                        voice.synthesize_wav(text, wav_file, syn_config=syn_config)
                    except Exception:
                        # Fallback for older piper versions
                        try:
                            voice.synthesize_wav(text, wav_file, length_scale=length_scale)
                        except TypeError:
                            voice.synthesize_wav(text, wav_file)

                if self._stop_flag:
                    return

                if volume < 0.99:
                    self._apply_volume(temp_path, volume)

                winsound.PlaySound(temp_path, winsound.SND_FILENAME | winsound.SND_ASYNC)

                try:
                    with wave.open(temp_path, "rb") as wf:
                        duration = wf.getnframes() / float(wf.getframerate())
                except Exception:
                    duration = max(1.5, len(text.split()) * 0.45)

                start = time.time()
                while not self._stop_flag and (time.time() - start) < duration + 0.4:
                    time.sleep(0.05)

                winsound.PlaySound(None, winsound.SND_PURGE)

            except Exception as e:
                print(f"Piper error: {e}")
            finally:
                self._is_playing = False
                if temp_path:
                    try:
                        os.unlink(temp_path)
                    except Exception:
                        pass

        threading.Thread(target=_run, daemon=True).start()

    def _apply_volume(self, wav_path: str, volume: float):
        import array

        with wave.open(wav_path, "rb") as wf:
            nchannels = wf.getnchannels()
            sampwidth = wf.getsampwidth()
            framerate = wf.getframerate()
            nframes = wf.getnframes()
            frames = wf.readframes(nframes)

        if sampwidth != 2:
            return

        samples = array.array("h")
        samples.frombytes(frames)

        for i in range(len(samples)):
            val = int(samples[i] * volume)
            if val > 32767:
                val = 32767
            elif val < -32768:
                val = -32768
            samples[i] = val

        with wave.open(wav_path, "wb") as wf:
            wf.setnchannels(nchannels)
            wf.setsampwidth(sampwidth)
            wf.setframerate(framerate)
            wf.writeframes(samples.tobytes())

    def is_speaking(self) -> bool:
        return self._is_playing

    def stop(self):
        self._stop_flag = True
        self._is_playing = False
        winsound.PlaySound(None, winsound.SND_PURGE)
# ---------------------------------------------------------------------------
# PowerPoint controller
# ---------------------------------------------------------------------------

class PPTController:
    """
    Manages PowerPoint note extraction (python-pptx) and COM slideshow control.

    IMPORTANT – COM threading rule:
      open_slideshow(), goto_slide(), and close() must ALL be called from the
      SAME thread (the presenter background thread).  COM objects are
      apartment-threaded and cannot be shared across threads.
    """

    def __init__(self, filepath: str):
        self.filepath = os.path.abspath(filepath)
        self.prs = Presentation(filepath)   # python-pptx – thread-safe, for notes
        self._app = None
        self._presentation = None
        self._slideshow = None

    # ---- note extraction (safe to call from any thread) -----------------

    def get_notes(self, slide_index: int) -> str:
        """
        Return the speaker notes for the given 0-based slide index.

        Uses direct XPath extraction instead of notes_text_frame so that it
        works regardless of placeholder index (which varies across PPTX files
        and PowerPoint versions).
        """
        slide = self.prs.slides[slide_index]
        notes_slide = slide.notes_slide
        if notes_slide is None:
            return ""
        texts = []
        for sp in notes_slide._element.findall('.//p:sp', _PPTX_NSMAP):
            ph = sp.find('.//p:ph', _PPTX_NSMAP)
            if ph is None:
                continue
            # Skip non-body placeholders (slide image, page number, date, etc.)
            if ph.get('type', 'body') in ('sldImg', 'sldNum', 'dt', 'hdr', 'ftr'):
                continue
            tx_body = sp.find('p:txBody', _PPTX_NSMAP)
            if tx_body is None:
                continue
            for t_elem in tx_body.findall('.//a:t', _PPTX_NSMAP):
                if t_elem.text:
                    texts.append(t_elem.text)
        return ' '.join(texts).strip()

    @property
    def slide_count(self) -> int:
        return len(self.prs.slides)

    # ---- COM slideshow control (must stay on ONE thread) ----------------

    def open_slideshow(self):
        """
        Launch PowerPoint and start the slideshow.
        Must be called on the presenter thread AFTER CoInitialize().
        """
        self._app = win32com.client.Dispatch("PowerPoint.Application")
        self._app.Visible = True
        self._presentation = self._app.Presentations.Open(
            self.filepath, ReadOnly=True, Untitled=False, WithWindow=True
        )
        self._presentation.SlideShowSettings.Run()
        # Wait for the SlideShowWindow to become available
        for _ in range(20):
            time.sleep(0.3)
            try:
                if self._presentation.SlideShowWindow is not None:
                    self._slideshow = self._presentation.SlideShowWindow.View
                    break
            except Exception:
                pass

    def goto_slide(self, slide_number: int):
        """Jump to a 1-based slide number. Must be called on the same thread as open_slideshow()."""
        if self._slideshow is None:
            return
        try:
            self._slideshow.GotoSlide(slide_number)
        except Exception:
            pass

    def close(self):
        """Close the presentation. Must be called on the same thread as open_slideshow()."""
        try:
            if self._presentation:
                self._presentation.Close()
        except Exception:
            pass
        try:
            if self._app:
                self._app.Quit()
        except Exception:
            pass
        finally:
            self._slideshow = None
            self._presentation = None
            self._app = None


# ---------------------------------------------------------------------------
# Presenter (orchestrates TTS + slide advancing)
# ---------------------------------------------------------------------------

class Presenter:
    """Runs the auto-presentation loop on a background thread."""

    def __init__(self, ppt: PPTController, tts: TTSEngine,
                 on_slide_change, on_status_change, on_finished,on_subtitle_update=None):
        self._ppt = ppt
        self._tts = tts
        self._on_slide_change = on_slide_change   # callback(slide_index)
        self._on_status_change = on_status_change # callback(message)
        self._on_finished = on_finished           # callback()
        self._on_subtitle_update = on_subtitle_update
        self._thread: threading.Thread | None = None
        self._stop_event = threading.Event()
        self._pause_event = threading.Event()
        self._pause_event.set()  # not paused initially
        self.current_slide = 0
        self._jump_to = None          # target 0-based slide index, or None
        self._jump_event = threading.Event()

    def next_slide(self):
        """Request jump to next slide (called from GUI thread)."""
        self._jump_to = self.current_slide + 1
        self._jump_event.set()
        self._tts.stop()
        self._pause_event.set()   # unblock if paused

    def prev_slide(self):
        """Request jump to previous slide (called from GUI thread)."""
        self._jump_to = max(0, self.current_slide - 1)
        self._jump_event.set()
        self._tts.stop()
        self._pause_event.set()

    def start(self, from_slide: int = 0):
        self._stop_event.clear()
        self._pause_event.set()
        self.current_slide = from_slide
        self._thread = threading.Thread(target=self._run, daemon=True)
        self._thread.start()

    def pause(self):
        self._pause_event.clear()
        # Do NOT call self._tts.stop() here – the presenter thread will do it

    def resume(self):
        self._pause_event.set()

    def stop(self):
        self._stop_event.set()
        self._pause_event.set()   # unblock any wait
        # again, do NOT call tts.stop() from the GUI thread
    def is_running(self) -> bool:
        return self._thread is not None and self._thread.is_alive()


    def _run(self):
        pythoncom.CoInitialize()
        try:
            if hasattr(self._tts, "_init_engine"):
                self._tts._init_engine()

            self._on_status_change("Opening PowerPoint slideshow…")
            try:
                self._ppt.open_slideshow()
            except Exception as e:
                self._on_status_change(f"Error opening slideshow: {e}")
                self._on_finished()
                return

            total = self._ppt.slide_count
            idx = self.current_slide

            while idx < total:
                if self._stop_event.is_set():
                    break

                self._pause_event.wait()
                if self._stop_event.is_set():
                    break

                if self._jump_event.is_set():
                    self._jump_event.clear()
                    if self._jump_to is not None and 0 <= self._jump_to < total:
                        idx = self._jump_to
                    self._jump_to = None

                self.current_slide = idx
                slide_num = idx + 1

                self._ppt.goto_slide(slide_num)
                self._on_slide_change(idx)
                self._on_status_change(f"Slide {slide_num} / {total}")

                notes = self._ppt.get_notes(idx)

                if notes:
                    chunks = split_into_chunks(notes)
                    chunk_idx = 0

                    while chunk_idx < len(chunks) and not self._stop_event.is_set():
                        if self._jump_event.is_set():
                            self._tts.stop()
                            break

                        self._on_status_change(
                            f"Slide {slide_num} / {total}  —  Speaking… ({chunk_idx+1}/{len(chunks)})"
                        )
                        # Update subtitle to the current sentence/chunk
                        if self._on_subtitle_update:
                            self._on_subtitle_update(chunks[chunk_idx])

                        self._tts.speak_async(chunks[chunk_idx])
                        time.sleep(0.4)

                        paused = False
                        while True:
                            if self._stop_event.is_set():
                                self._tts.stop()
                                break
                            if self._jump_event.is_set():
                                self._tts.stop()
                                break
                            if not self._pause_event.is_set():
                                self._tts.stop()
                                paused = True
                                break
                            if not self._tts.is_speaking():
                                break
                            time.sleep(0.08)

                        if self._stop_event.is_set() or self._jump_event.is_set():
                            break

                        if paused:
                            while not self._pause_event.is_set():
                                if self._stop_event.is_set() or self._jump_event.is_set():
                                    break
                                time.sleep(0.05)
                            if self._stop_event.is_set() or self._jump_event.is_set():
                                break
                            continue

                        chunk_idx += 1
                else:
                    self._on_status_change(
                        f"Slide {slide_num} / {total}  —  No notes, waiting 3 s…"
                    )
                    if self._on_subtitle_update:
                            self._on_subtitle_update("(No notes)")
                    for _ in range(30):
                        if self._stop_event.is_set() or self._jump_event.is_set():
                            break
                        self._pause_event.wait()
                        if self._stop_event.is_set() or self._jump_event.is_set():
                            break
                        time.sleep(0.1)

                if self._stop_event.is_set():
                    break

                if self._jump_event.is_set():
                    continue

                idx += 1

            if not self._stop_event.is_set():
                self._on_status_change("Presentation finished.")
                self._on_finished()

        finally:
            self._ppt.close()
            pythoncom.CoUninitialize()

# ---------------------------------------------------------------------------
# GUI
# ---------------------------------------------------------------------------
class SubtitleWindow(tk.Toplevel):
    """subtitle window – appears at the bottom of the screen."""

    def __init__(self, master):
        super().__init__(master)
        self.title("Subtitles")
        self.configure(bg="#000000")
        self.attributes("-topmost", True)      # always on top
        self.attributes("-alpha", 0.88)        # slightly transparent
        self.resizable(True, True)

        # Get screen size
        screen_width = self.winfo_screenwidth()
        screen_height = self.winfo_screenheight()

        # Window size: wide but short
        win_width = min(1100, screen_width - 100)
        win_height = 130

        # Position at bottom center
        x = (screen_width - win_width) // 2
        y = screen_height - win_height - 40     # 40px from bottom

        self.geometry(f"{win_width}x{win_height}+{x}+{y}")

        self._label = tk.Label(
            self,
            text="(No notes)",
            font=("Segoe UI", 20),
            fg="white",
            bg="#000000",
            wraplength=win_width - 40,
            justify="center",
            anchor="center"
        )
        self._label.pack(expand=True, fill="both", padx=15, pady=10)

        # Handle close button – just hide the window
        self.protocol("WM_DELETE_WINDOW", self._on_close)

    def update_text(self, text: str):
        if not text or not text.strip():
            self._label.config(text="(No notes)")
        else:
            self._label.config(text=text.strip())

    def _on_close(self):
        self.withdraw()

class AutoPresentApp(tk.Tk):
    def __init__(self):
        super().__init__()
        self.title("AutoPresent")
        self.resizable(False, False)
        self.configure(bg="#1e1e2e")

        self._ppt: PPTController | None = None
        self._tts_sapi = TTSEngine()
        try:
            self._tts_piper = PiperTTSEngine()
            self._has_piper = True
        except Exception as e:
            print(f"Piper not available: {e}")
            self._tts_piper = None
            self._has_piper = False

        self._tts = self._tts_sapi  # default
        self._presenter: Presenter | None = None
        self._paused = False
        self._subtitle_win: SubtitleWindow | None = None
        self._show_subtitles = tk.BooleanVar(value=False)

        self._build_ui()
        self.protocol("WM_DELETE_WINDOW", self._on_close)

    # ---- UI construction -------------------------------------------------

    def _build_ui(self):
        PAD = 12
        BG = "#1e1e2e"
        FG = "#cdd6f4"
        ACCENT = "#89b4fa"
        ENTRY_BG = "#313244"
        BTN_BG = "#45475a"
        BTN_FG = "#cdd6f4"

        style = ttk.Style(self)
        style.theme_use("clam")
        style.configure("TScale", background=BG, troughcolor=ENTRY_BG,
                        sliderlength=18, sliderrelief="flat")
        style.configure("TCombobox", fieldbackground=ENTRY_BG,
                        background=BTN_BG, foreground=FG,
                        selectbackground=ENTRY_BG, selectforeground=FG)
        style.map("TCombobox", fieldbackground=[("readonly", ENTRY_BG)])

        # ---- File row ----
        file_frame = tk.Frame(self, bg=BG)
        file_frame.grid(row=0, column=0, columnspan=3,
                        padx=PAD, pady=(PAD, 4), sticky="ew")

        tk.Label(file_frame, text="PowerPoint file:", bg=BG, fg=FG,
                 font=("Segoe UI", 10)).pack(side="left")

        self._file_var = tk.StringVar(value="(none selected)")
        tk.Label(file_frame, textvariable=self._file_var, bg=BG,
                 fg=ACCENT, font=("Segoe UI", 10, "italic"),
                 width=40, anchor="w").pack(side="left", padx=8)

        tk.Button(file_frame, text="Browse…", bg=BTN_BG, fg=BTN_FG,
                  font=("Segoe UI", 9), relief="flat", padx=6,
                  command=self._browse_file).pack(side="left")

        # ---- Settings ----
        settings_frame = tk.LabelFrame(self, text=" Settings ",
                                       bg=BG, fg=FG,
                                       font=("Segoe UI", 9, "bold"),
                                       labelanchor="nw", bd=1,
                                       relief="groove")
        settings_frame.grid(row=1, column=0, columnspan=3,
                            padx=PAD, pady=4, sticky="ew")
        settings_frame.columnconfigure(1, weight=1)

        # Speech rate
        tk.Label(settings_frame, text="Speech rate:", bg=BG, fg=FG,
                 font=("Segoe UI", 9)).grid(row=0, column=0,
                                            padx=8, pady=4, sticky="w")
        self._rate_var = tk.IntVar(value=175)
        rate_scale = ttk.Scale(settings_frame, from_=80, to=300,
                               variable=self._rate_var, orient="horizontal",
                               length=200,
                               command=lambda v: self._on_rate_change())
        rate_scale.grid(row=0, column=1, padx=8, pady=4, sticky="ew")
        self._rate_label = tk.Label(settings_frame, text="175 wpm",
                                    bg=BG, fg=FG, font=("Segoe UI", 9),
                                    width=8)
        self._rate_label.grid(row=0, column=2, padx=(0, 8))

        # Volume
        tk.Label(settings_frame, text="Volume:", bg=BG, fg=FG,
                 font=("Segoe UI", 9)).grid(row=1, column=0,
                                            padx=8, pady=4, sticky="w")
        self._vol_var = tk.DoubleVar(value=1.0)
        vol_scale = ttk.Scale(settings_frame, from_=0.0, to=1.0,
                              variable=self._vol_var, orient="horizontal",
                              length=200,
                              command=lambda v: self._on_vol_change())
        vol_scale.grid(row=1, column=1, padx=8, pady=4, sticky="ew")
        self._vol_label = tk.Label(settings_frame, text="100%",
                                   bg=BG, fg=FG, font=("Segoe UI", 9),
                                   width=8)
        self._vol_label.grid(row=1, column=2, padx=(0, 8))

        # Engine selector
        tk.Label(settings_frame, text="Engine:", bg=BG, fg=FG,
                 font=("Segoe UI", 9)).grid(row=2, column=0, padx=8, pady=4, sticky="w")

        engine_values = ["SAPI (Windows)"]
        if getattr(self, "_has_piper", False):
            engine_values.append("Piper (Neural)")

        self._engine_var = tk.StringVar(value="SAPI (Windows)")
        self._engine_combo = ttk.Combobox(settings_frame,
                                          textvariable=self._engine_var,
                                          values=engine_values,
                                          state="readonly", width=35)
        self._engine_combo.grid(row=2, column=1, columnspan=2, padx=8, pady=4, sticky="ew")
        self._engine_combo.bind("<<ComboboxSelected>>", self._on_engine_change)

        # Voice selector
        tk.Label(settings_frame, text="Voice:", bg=BG, fg=FG,
                 font=("Segoe UI", 9)).grid(row=3, column=0,
                                            padx=8, pady=4, sticky="w")
        self._voice_var = tk.StringVar()
        self._voice_combo = ttk.Combobox(settings_frame,
                                         textvariable=self._voice_var,
                                         state="readonly", width=35)
        self._voice_combo.grid(row=3, column=1, columnspan=2,
                               padx=8, pady=4, sticky="ew")
        self._populate_voices()
        self._voice_combo.bind("<<ComboboxSelected>>", self._on_voice_change)

        # Start from slide + Show Subtitles (closer together)
        tk.Label(settings_frame, text="Start from slide:", bg=BG, fg=FG,
                font=("Segoe UI", 9)).grid(row=4, column=0,
                                            padx=8, pady=4, sticky="w")

        # Frame to hold spinbox + checkbox side by side
        start_frame = tk.Frame(settings_frame, bg=BG)
        start_frame.grid(row=4, column=1, columnspan=2, padx=8, pady=4, sticky="w")

        self._start_slide_var = tk.IntVar(value=1)
        self._start_slide_spin = tk.Spinbox(
            start_frame, from_=1, to=999,
            textvariable=self._start_slide_var,
            width=6, bg=ENTRY_BG, fg=FG,
            buttonbackground=BTN_BG, relief="flat",
            font=("Segoe UI", 9)
        )
        self._start_slide_spin.pack(side="left")

        tk.Checkbutton(
            start_frame,
            text="Show Subtitles window",
            variable=self._show_subtitles,
            bg=BG,
            fg=FG,
            selectcolor=ENTRY_BG,
            activebackground=BG,
            activeforeground=FG,
            font=("Segoe UI", 9),
            command=self._toggle_subtitles
        ).pack(side="left", padx=(12, 0))
        # ---- Progress / status ----
        prog_frame = tk.Frame(self, bg=BG)
        prog_frame.grid(row=2, column=0, columnspan=3,
                        padx=PAD, pady=4, sticky="ew")
        prog_frame.columnconfigure(0, weight=1)

        self._status_var = tk.StringVar(value="Ready. Open a .pptx file to begin.")
        tk.Label(prog_frame, textvariable=self._status_var,
                 bg=BG, fg=FG, font=("Segoe UI", 9),
                 anchor="w").grid(row=0, column=0, sticky="ew")

        self._progress = ttk.Progressbar(prog_frame, orient="horizontal",
                                         mode="determinate", length=460)
        self._progress.grid(row=1, column=0, pady=(2, 0), sticky="ew")

        # Current slide notes preview
        self._notes_text = tk.Text(self, height=7, width=58, wrap="word",
                                   bg=ENTRY_BG, fg=FG, relief="flat",
                                   font=("Segoe UI", 9),
                                   state="disabled", bd=0,
                                   padx=6, pady=4)
        self._notes_text.grid(row=3, column=0, columnspan=3,
                              padx=PAD, pady=4)

        # ---- Control buttons ----
        btn_frame = tk.Frame(self, bg=BG)
        btn_frame.grid(row=4, column=0, columnspan=3,
                       padx=PAD, pady=(4, PAD))

        btn_cfg = dict(font=("Segoe UI", 10, "bold"), relief="flat",
                       width=12, padx=4, pady=6)

        self._start_btn = tk.Button(btn_frame, text="▶  Start",
                                    bg="#a6e3a1", fg="#1e1e2e",
                                    command=self._on_start, **btn_cfg)
        self._start_btn.pack(side="left", padx=4)

        self._pause_btn = tk.Button(btn_frame, text="⏸  Pause",
                                    bg="#f9e2af", fg="#1e1e2e",
                                    command=self._on_pause_resume,
                                    state="disabled", **btn_cfg)
        self._pause_btn.pack(side="left", padx=4)

        self._stop_btn = tk.Button(btn_frame, text="⏹  Stop",
                                   bg="#f38ba8", fg="#1e1e2e",
                                   command=self._on_stop,
                                   state="disabled", **btn_cfg)
        self._stop_btn.pack(side="left", padx=4)

        self._prev_btn = tk.Button(btn_frame, text="⏮  Prev",
                                   bg=BTN_BG, fg=BTN_FG,
                                   command=self._on_prev,
                                   state="disabled", **btn_cfg)
        self._prev_btn.pack(side="left", padx=4)

        self._next_btn = tk.Button(btn_frame, text="Next  ⏭",
                                   bg=BTN_BG, fg=BTN_FG,
                                   command=self._on_next,
                                   state="disabled", **btn_cfg)
        self._next_btn.pack(side="left", padx=4)


    # ---- Callbacks -------------------------------------------------------

    def _browse_file(self):
        path = filedialog.askopenfilename(
            title="Select PowerPoint file",
            filetypes=[("PowerPoint files", "*.pptx *.ppt"), ("All files", "*.*")]
        )
        if not path:
            return
        self._load_file(path)

    def _load_file(self, path: str):
        try:
            self._ppt = PPTController(path)
            short = os.path.basename(path)
            self._file_var.set(short)
            n = self._ppt.slide_count
            self._progress["maximum"] = n
            self._progress["value"] = 0
            self._start_slide_spin.config(to=n)
            self._start_slide_var.set(1)
            self._status_var.set(f"Loaded: {short}  ({n} slides)")
            self._show_notes(0)
            self._start_btn.config(state="normal")
        except Exception as e:
            messagebox.showerror("Error", f"Could not open file:\n{e}")

    def _on_rate_change(self):
        rate = self._rate_var.get()
        self._rate_label.config(text=f"{rate} wpm")
        self._tts.set_rate(rate)

    def _on_vol_change(self):
        vol = self._vol_var.get()
        self._vol_label.config(text=f"{int(vol * 100)}%")
        self._tts.set_volume(vol)

    def _on_voice_change(self, _event=None):
        idx = self._voice_combo.current()
        if idx >= 0 and idx < len(self._voices):
            self._tts.set_voice(self._voices[idx].id)

    def _on_engine_change(self, event=None):
        engine = self._engine_var.get()
        if engine.startswith("Piper") and self._has_piper:
            self._tts = self._tts_piper
        else:
            self._tts = self._tts_sapi
        self._populate_voices()

    def _populate_voices(self):
        try:
            voices = self._tts.get_voices_sync()
            names = [v.name for v in voices]
            self._voices = voices
            self._voice_combo["values"] = names
            if names:
                self._voice_combo.current(0)
                self._voice_var.set(names[0])
                self._tts.set_voice(voices[0].id)
        except Exception as e:
            messagebox.showerror("Error", f"Could not load voices:\n{e}")

    def _on_start(self):
        if self._ppt is None:
            messagebox.showwarning("No file", "Please select a .pptx file first.")
            return

        # Apply current settings
        self._tts.set_rate(self._rate_var.get())
        self._tts.set_volume(self._vol_var.get())
        self._on_voice_change()

        start_idx = self._start_slide_var.get() - 1  # convert to 0-based

        self._paused = False
        self._presenter = Presenter(
            ppt=self._ppt,
            tts=self._tts,
            on_slide_change=self._cb_slide_change,
            on_status_change=self._cb_status,
            on_finished=self._cb_finished,
            on_subtitle_update=self._cb_subtitle_update,
        )
        self._presenter.start(from_slide=start_idx)

        self._start_btn.config(state="disabled")
        self._pause_btn.config(state="normal")
        self._stop_btn.config(state="normal")
        self._prev_btn.config(state="normal")
        self._next_btn.config(state="normal")

    def _on_pause_resume(self):
        if self._presenter is None:
            return
        if self._paused:
            self._presenter.resume()
            self._pause_btn.config(text="⏸  Pause")
            self._paused = False
        else:
            self._presenter.pause()
            self._pause_btn.config(text="▶  Resume")
            self._paused = True

    def _on_stop(self):
        if self._presenter:
            # stop() signals the presenter thread; it will call ppt.close()
            # on its own thread before exiting, so we don't close here.
            self._presenter.stop()
        self._reset_controls()
        self._status_var.set("Stopped.")

    def _on_prev(self):
        if self._presenter and self._presenter.is_running():
            self._presenter.prev_slide()

    def _on_next(self):
        if self._presenter and self._presenter.is_running():
            self._presenter.next_slide()

    def _reset_controls(self):
        self._start_btn.config(state="normal" if self._ppt else "disabled")
        self._pause_btn.config(state="disabled", text="⏸  Pause")
        self._stop_btn.config(state="disabled")
        self._prev_btn.config(state="disabled")
        self._next_btn.config(state="disabled")
        self._paused = False

    # ---- Thread-safe UI callbacks ----------------------------------------

    def _cb_slide_change(self, slide_idx: int):
        self.after(0, self._update_slide_ui, slide_idx)

    def _cb_status(self, message: str):
        self.after(0, self._status_var.set, message)

    def _cb_finished(self):
        self.after(0, self._on_presentation_finished)

    def _cb_subtitle_update(self, text: str):
        self.after(0, self._update_subtitles, text)

    def _update_slide_ui(self, slide_idx: int):
        self._progress["value"] = slide_idx + 1
        self._show_notes(slide_idx)

    def _show_notes(self, slide_idx: int):
        if self._ppt is None:
            return
        notes = self._ppt.get_notes(slide_idx)
        # Update the main notes box
        self._notes_text.config(state="normal")
        self._notes_text.delete("1.0", "end")
        if notes:
            self._notes_text.insert("end", notes)
        else:
            self._notes_text.insert("end", "(No notes for this slide)")
        self._notes_text.config(state="disabled")
        # Update subtitle window
        self._update_subtitles(notes)

    def _toggle_subtitles(self):
        if self._show_subtitles.get():
            if self._subtitle_win is None or not self._subtitle_win.winfo_exists():
                self._subtitle_win = SubtitleWindow(self)
            else:
                self._subtitle_win.deiconify()
                self._subtitle_win.attributes("-topmost", True)

            # Show current notes immediately
            if self._ppt is not None:
                current_idx = max(0, self._progress["value"] - 1)
                notes = self._ppt.get_notes(int(current_idx))
                self._subtitle_win.update_text(notes)
        else:
            if self._subtitle_win is not None and self._subtitle_win.winfo_exists():
                self._subtitle_win.withdraw()

    def _update_subtitles(self, text: str):
        if self._subtitle_win is not None and self._subtitle_win.winfo_exists():
            if self._show_subtitles.get():
                self._subtitle_win.update_text(text)

    def _on_presentation_finished(self):
        self._reset_controls()
        self._progress["value"] = self._ppt.slide_count if self._ppt else 0
        messagebox.showinfo("AutoPresent", "Presentation finished!")

    # ---- Window close ----------------------------------------------------

    def _on_close(self):
        if self._subtitle_win is not None and self._subtitle_win.winfo_exists():
            self._subtitle_win.destroy()

        if self._presenter and self._presenter.is_running():
            if not messagebox.askyesno("Quit",
                                       "Presentation is running. Quit anyway?"):
                return
            self._presenter.stop()
            # Give the presenter thread a moment to close COM cleanly
            self._presenter._thread.join(timeout=3)
        self.destroy()


# ---------------------------------------------------------------------------
# Entry point
# ---------------------------------------------------------------------------

if __name__ == "__main__":
    app = AutoPresentApp()
    app.mainloop()
